#!/usr/bin/env python3
"""Generate Mealvana Endurance pitch deck with brand styling."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Brand colors ──
BLACKBERRY = RGBColor(0x3D, 0x1F, 0x47)
BLACKBERRY_LIGHT = RGBColor(0x4A, 0x28, 0x54)
ORANGE = RGBColor(0xFF, 0x8B, 0x3D)
ELECTROLYTE = RGBColor(0x5D, 0xE4, 0xD3)
CREAM = RGBColor(0xF5, 0xF3, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x38, 0x16, 0x33)
GRAY = RGBColor(0x99, 0x99, 0x99)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DRAGONFRUIT = RGBColor(0xE8, 0x43, 0x93)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, opacity=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if opacity is not None:
        # Set fill transparency via XML
        solidFill = shape.fill._fill
        srgb = solidFill.find(qn('a:solidFill')).find(qn('a:srgbClr'))
        if srgb is not None:
            alpha = srgb.makeelement(qn('a:alpha'), {})
            alpha.set('val', str(int((1 - opacity) * 100000)))
            srgb.append(alpha)
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri', line_spacing=1.2):
    """Add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=18,
                   default_color=WHITE, alignment=PP_ALIGN.LEFT):
    """Add text box with multiple styled paragraphs.
    lines: list of (text, font_size, color, bold) tuples
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Calibri'
        p.alignment = alignment
        p.space_after = Pt(4)
    return txBox


def add_accent_bar(slide, left, top, width=Inches(0.6), height=Pt(4), color=ORANGE):
    """Add a small accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_slide_number(slide, num, total=12):
    """Add subtle slide number."""
    add_text_box(slide, Inches(12.4), Inches(7.0), Inches(0.8), Inches(0.4),
                 f'{num}', font_size=10, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════
# SLIDE 1: COVER
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BLACKBERRY)

# Accent stripe at top
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)

# Company name + tagline
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
             'MEALVANA ENDURANCE', font_size=16, color=ELECTROLYTE,
             bold=False, alignment=PP_ALIGN.LEFT)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(2),
             'The Nutrition Intelligence\nLayer for Endurance Sports',
             font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

add_accent_bar(slide, Inches(1.5), Inches(4.8), Inches(1.2), Pt(4), ORANGE)

add_text_box(slide, Inches(1.5), Inches(5.1), Inches(10), Inches(0.8),
             'AI-powered, personalized nutrition planning for runners, cyclists, and triathletes.',
             font_size=20, color=CREAM, bold=False, alignment=PP_ALIGN.LEFT)

add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
             'Xuan Huang, CEO  |  Lee Martin, CTO  |  April 2026',
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)


# ════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(3), Inches(0.5),
             'THE PROBLEM', font_size=14, color=ORANGE, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.1), Inches(10), Inches(1.2),
             'Endurance athletes don\'t know what to eat,\nwhen to eat it, or how much.',
             font_size=36, color=BLACKBERRY, bold=True)

add_accent_bar(slide, Inches(1.2), Inches(2.5), Inches(0.8), Pt(3), ORANGE)

# Stat boxes
stats = [
    ('30-50%', 'of marathon runners\nexperience GI distress', 'ACSM'),
    ('10-20%', 'performance loss from\ninadequate fueling', 'Research'),
    ('$100-300+', 'wasted per race from\nnutrition mistakes', 'Entry fees + travel'),
    ('0', 'personalized race-day\nnutrition tools exist', 'Market gap'),
]

for i, (big_num, desc, source) in enumerate(stats):
    x = Inches(1.2 + i * 2.9)
    y = Inches(3.2)

    # Card background
    card = add_shape(slide, x, y, Inches(2.6), Inches(3.2), WHITE)
    card.shadow.inherit = False

    add_text_box(slide, x + Inches(0.3), y + Inches(0.4), Inches(2), Inches(0.8),
                 big_num, font_size=36, color=ORANGE, bold=True)

    add_text_box(slide, x + Inches(0.3), y + Inches(1.4), Inches(2), Inches(1.0),
                 desc, font_size=16, color=TEXT_DARK, bold=False)

    add_text_box(slide, x + Inches(0.3), y + Inches(2.5), Inches(2), Inches(0.4),
                 source, font_size=11, color=GRAY, bold=False)

add_slide_number(slide, 2)


# ════════════════════════════════════════════
# SLIDE 3: THE SOLUTION
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(3), Inches(0.5),
             'THE SOLUTION', font_size=14, color=ORANGE, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.1), Inches(10), Inches(1.2),
             'Personalized, science-based nutrition plans\ngenerated by AI in seconds.',
             font_size=36, color=WHITE, bold=True)

add_accent_bar(slide, Inches(1.2), Inches(2.5), Inches(0.8), Pt(3), ELECTROLYTE)

# Feature pillars
pillars = [
    ('For Athletes', [
        'AI-generated race-day fueling plans',
        'Pre / during / post workout phases',
        'Food preference learning (love / try / avoid)',
        'Weather-adjusted hydration targets',
        'Carb-loading protocols',
    ]),
    ('For Coaches', [
        'Web dashboard + mobile app',
        'Many-to-many athlete management',
        'Edit & monitor nutrition plans',
        'In-app messaging',
        'Certification verification',
    ]),
    ('Science-Backed', [
        'ACSM metabolic formulas',
        'ISSN sports nutrition guidelines',
        'Dual AI + algorithmic engine',
        'Clinical review by sports dietitian',
        'Offline-first reliability',
    ]),
]

for i, (title, items) in enumerate(pillars):
    x = Inches(1.2 + i * 3.8)
    y = Inches(3.2)

    add_text_box(slide, x, y, Inches(3.4), Inches(0.5),
                 title, font_size=20, color=ELECTROLYTE, bold=True)

    for j, item in enumerate(items):
        add_text_box(slide, x, y + Inches(0.6 + j * 0.55), Inches(3.4), Inches(0.5),
                     f'\u2022  {item}', font_size=14, color=CREAM)

add_slide_number(slide, 3)


# ════════════════════════════════════════════
# SLIDE 4: PRODUCT
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(3), Inches(0.5),
             'THE PRODUCT', font_size=14, color=ORANGE, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.1), Inches(10), Inches(1.0),
             'A complete nutrition planning platform.',
             font_size=36, color=BLACKBERRY, bold=True)

# Feature highlight cards
features = [
    ('AI Plan Generation', 'Personalized nutrition plans in seconds using ACSM formulas + LLM intelligence'),
    ('Food Preferences', '3-tier system learns what athletes love, will try, and want to avoid'),
    ('Coach Mode', 'Web dashboard for coaches to manage athletes and edit plans at scale'),
    ('Integrations', 'TrainingPeaks, FinalSurge, Garmin \u2014 sync activities and push nutrition data back'),
    ('Carb Loading', 'Multi-day protocols tailored to race distance and athlete weight'),
    ('Offline-First', 'Full functionality without internet \u2014 critical for remote training'),
]

for i, (title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(1.2 + col * 3.8)
    y = Inches(2.4 + row * 2.3)

    card = add_shape(slide, x, y, Inches(3.4), Inches(2.0), WHITE)

    add_accent_bar(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.5), Pt(3), ORANGE)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.5), Inches(2.8), Inches(0.5),
                 title, font_size=18, color=BLACKBERRY, bold=True)

    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(2.8), Inches(0.8),
                 desc, font_size=13, color=TEXT_DARK)

# Placeholder note for screenshots
add_text_box(slide, Inches(8.5), Inches(6.5), Inches(4.5), Inches(0.5),
             '[Add app screenshots here]', font_size=12, color=GRAY,
             alignment=PP_ALIGN.RIGHT)

add_slide_number(slide, 4)


# ════════════════════════════════════════════
# SLIDE 5: WHY NOW
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(3), Inches(0.5),
             'WHY NOW', font_size=14, color=ORANGE, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.1), Inches(10), Inches(1.0),
             'Five converging trends.',
             font_size=36, color=WHITE, bold=True)

trends = [
    ('Record participation', '650K+ marathon finishers in 2023.\n2M+ half-marathon participants.\nTriathlon growing 8-10% annually.'),
    ('AI personalization wave', 'Consumer expectation has shifted\nfrom generic to individually tailored.\nLLMs make this possible at scale.'),
    ('Coaching going digital', 'Training is digital (TrainingPeaks,\nStrava). Nutrition is the last pillar\nstill done with spreadsheets.'),
    ('Nutrition brands booming', 'Gels, bars, hydration brands need\ndata-driven customer channels.\nOur recommendations = distribution.'),
    ('Remote coach economy', 'COVID accelerated remote coaching.\nCoaches need digital nutrition tools\nto serve athletes at scale.'),
]

for i, (title, desc) in enumerate(trends):
    x = Inches(0.8 + i * 2.5)
    y = Inches(2.5)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.05), y, Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    add_text_box(slide, x, y + Inches(0.8), Inches(2.3), Inches(0.5),
                 title, font_size=16, color=ELECTROLYTE, bold=True)

    add_text_box(slide, x, y + Inches(1.4), Inches(2.3), Inches(2.5),
                 desc, font_size=13, color=CREAM, line_spacing=1.4)

add_slide_number(slide, 5)


# ════════════════════════════════════════════
# SLIDE 6: MARKET SIZE
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(3), Inches(0.5),
             'MARKET SIZE', font_size=14, color=ORANGE, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.1), Inches(10), Inches(1.0),
             'Large, growing, underserved.',
             font_size=36, color=BLACKBERRY, bold=True)

# TAM / SAM / SOM circles (nested visual)
market_data = [
    ('TAM', '$4.2B', 'Digital sports nutrition\n& fitness apps globally', Inches(5.0), BLACKBERRY_LIGHT),
    ('SAM', '$680M', 'Endurance athletes in\nEnglish-speaking markets', Inches(3.6), BLACKBERRY),
    ('SOM', '$2-5M', '5-year capture:\n10K-25K subscribers', Inches(2.2), ORANGE),
]

center_x = Inches(3.5)
center_y = Inches(4.5)

for label, value, desc, diameter, color in market_data:
    x = center_x - diameter / 2 + Inches(0.5)
    y = center_y - diameter / 2
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, diameter, diameter)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

# Labels for circles
add_text_box(slide, Inches(1.2), Inches(2.5), Inches(1.5), Inches(0.4),
             'TAM', font_size=14, color=CREAM, bold=True)
add_text_box(slide, Inches(1.2), Inches(2.9), Inches(1.5), Inches(0.4),
             '$4.2B', font_size=28, color=CREAM, bold=True)

add_text_box(slide, Inches(2.5), Inches(3.8), Inches(2), Inches(0.4),
             'SAM  $680M', font_size=18, color=CREAM, bold=True)

add_text_box(slide, Inches(3.2), Inches(4.3), Inches(1.8), Inches(0.4),
             'SOM  $2-5M', font_size=16, color=WHITE, bold=True)

# Methodology on the right
add_text_box(slide, Inches(7.5), Inches(2.5), Inches(5), Inches(0.5),
             'Bottom-Up Methodology', font_size=18, color=BLACKBERRY, bold=True)

method_lines = [
    '60M active runners in the US (Running USA)',
    '17M competitive / event participants',
    '1-3% free-to-paid conversion (industry benchmark)',
    '10M+ US race registrations/year (RunSignUp)',
    'TrainingPeaks: 1M+ users, Strava: 120M+',
    '',
    'Year 3 target: $470K ARR',
    '  5,000 B2C subscribers @ $70/yr',
    '  200 coaches @ $240-600/yr',
    '',
    'Year 5 target: $2M-5M ARR',
    '  Geographic expansion + integration partnerships',
]

for i, line in enumerate(method_lines):
    if line == '':
        continue
    y_pos = Inches(3.1 + i * 0.32)
    is_bold = line.startswith('Year')
    color = ORANGE if is_bold else TEXT_DARK
    add_text_box(slide, Inches(7.5), y_pos, Inches(5), Inches(0.35),
                 line, font_size=13, color=color, bold=is_bold)

add_slide_number(slide, 6)


# ════════════════════════════════════════════
# SLIDE 7: TRACTION
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(3), Inches(0.5),
             'TRACTION', font_size=14, color=ORANGE, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.1), Inches(10), Inches(1.0),
             'Pre-revenue with strong product signals.',
             font_size=36, color=WHITE, bold=True)

# Traction metrics in a grid
traction = [
    ('100-500', 'Organic Users', 'Zero paid acquisition'),
    ('18', 'Versions Shipped', 'Rapid iteration since v1.0'),
    ('3', 'Live Integrations', 'TrainingPeaks, FinalSurge, Garmin'),
    ('$25K', 'AL Launchpad Prize', 'External concept validation'),
    ('11', 'Pro Features Gated', 'RevenueCat nearly complete'),
    ('2', 'Platforms', 'iOS + Android + Web portal'),
]

for i, (metric, label, detail) in enumerate(traction):
    col = i % 3
    row = i // 3
    x = Inches(1.2 + col * 3.8)
    y = Inches(2.5 + row * 2.4)

    # Card
    card = add_shape(slide, x, y, Inches(3.4), Inches(2.0), BLACKBERRY_LIGHT)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.8), Inches(0.8),
                 metric, font_size=40, color=ORANGE, bold=True)

    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(2.8), Inches(0.4),
                 label, font_size=18, color=WHITE, bold=True)

    add_text_box(slide, x + Inches(0.3), y + Inches(1.4), Inches(2.8), Inches(0.4),
                 detail, font_size=13, color=LIGHT_GRAY)

add_slide_number(slide, 7)


# ════════════════════════════════════════════
# SLIDE 8: BUSINESS MODEL
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(3), Inches(0.5),
             'BUSINESS MODEL', font_size=14, color=ORANGE, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.1), Inches(10), Inches(1.0),
             'Dual revenue: B2C subscriptions + B2B coach platform.',
             font_size=32, color=BLACKBERRY, bold=True)

# B2C column
add_shape(slide, Inches(1.2), Inches(2.5), Inches(5.2), Inches(4.5), WHITE)
add_accent_bar(slide, Inches(1.5), Inches(2.8), Inches(0.6), Pt(3), ORANGE)
add_text_box(slide, Inches(1.5), Inches(3.0), Inches(4.5), Inches(0.5),
             'B2C \u2014 Individual Athletes', font_size=22, color=BLACKBERRY, bold=True)

b2c_lines = [
    'Free tier: Basic plan generation + food preferences',
    'Pro tier: $9.99/mo  or  $69.99/yr',
    '',
    'Pro features include:',
    '  Training platform integrations',
    '  Coach Mode access',
    '  Brick workout nutrition',
    '  Carb-loading protocols',
    '  Personal fueling templates',
    '  Barcode scanning',
    '',
    'Payment via RevenueCat (Apple / Google)',
]

for i, line in enumerate(b2c_lines):
    if line == '':
        continue
    is_price = '$9.99' in line
    color = ORANGE if is_price else TEXT_DARK
    bold = is_price
    add_text_box(slide, Inches(1.5), Inches(3.6 + i * 0.3), Inches(4.5), Inches(0.3),
                 line, font_size=13, color=color, bold=bold)

# B2B column
add_shape(slide, Inches(7.0), Inches(2.5), Inches(5.2), Inches(4.5), WHITE)
add_accent_bar(slide, Inches(7.3), Inches(2.8), Inches(0.6), Pt(3), ELECTROLYTE)
add_text_box(slide, Inches(7.3), Inches(3.0), Inches(4.5), Inches(0.5),
             'B2B \u2014 Coaches & Dietitians', font_size=22, color=BLACKBERRY, bold=True)

b2b_lines = [
    'Starter: $20/mo (up to 20 athletes)',
    'Pro: $50/mo (unlimited athletes)',
    '',
    'Coach Mode includes:',
    '  Web dashboard + mobile',
    '  Many-to-many athlete management',
    '  Plan editing & monitoring',
    '  In-app messaging',
    '  Certification verification',
    '',
    'Sales: Direct to RRCA / USAT / NASM coaches',
    'Payment via Stripe',
]

for i, line in enumerate(b2b_lines):
    if line == '':
        continue
    is_price = '$20' in line or '$50' in line
    color = ORANGE if is_price else TEXT_DARK
    bold = is_price
    add_text_box(slide, Inches(7.3), Inches(3.6 + i * 0.3), Inches(4.5), Inches(0.3),
                 line, font_size=13, color=color, bold=bold)

add_slide_number(slide, 8)


# ════════════════════════════════════════════
# SLIDE 9: GO-TO-MARKET
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(3), Inches(0.5),
             'GO-TO-MARKET', font_size=14, color=ORANGE, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.1), Inches(10), Inches(1.0),
             'Three channels to first 1,000 users.',
             font_size=36, color=WHITE, bold=True)

channels = [
    ('1', 'App Store + Paid Acquisition',
     'Apple Search Ads and Google Ads targeting endurance keywords (marathon nutrition, race fueling, triathlon nutrition plan). Establish cost-per-install baseline.',
     'Target: 500+ installs'),
    ('2', 'Coach Network Effect',
     'Onboard 10+ certified coaches (RRCA, USAT). Each coach brings 10-50 athletes. Built-in viral distribution through coach-athlete relationships.',
     'Target: 200+ athletes via coaches'),
    ('3', 'Event Partnerships',
     'Partner with regional race organizers (Mercedes Marathon, Vulcan Run) to offer Mealvana as value-add to registered participants. Batch user acquisition.',
     'Target: 300+ event users'),
]

for i, (num, title, desc, target) in enumerate(channels):
    x = Inches(1.2 + i * 3.8)
    y = Inches(2.5)

    card = add_shape(slide, x, y, Inches(3.4), Inches(4.3), BLACKBERRY_LIGHT)

    # Number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.3), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(2.8), Inches(0.7),
                 title, font_size=18, color=ELECTROLYTE, bold=True)

    add_text_box(slide, x + Inches(0.3), y + Inches(1.7), Inches(2.8), Inches(1.8),
                 desc, font_size=13, color=CREAM, line_spacing=1.4)

    add_text_box(slide, x + Inches(0.3), y + Inches(3.6), Inches(2.8), Inches(0.4),
                 target, font_size=14, color=ORANGE, bold=True)

add_slide_number(slide, 9)


# ════════════════════════════════════════════
# SLIDE 10: COMPETITIVE LANDSCAPE
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(3), Inches(0.5),
             'COMPETITION', font_size=14, color=ORANGE, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.1), Inches(10), Inches(1.0),
             'We win on depth, personalization, and the coach platform.',
             font_size=32, color=BLACKBERRY, bold=True)

# Comparison table header
table_x = Inches(1.2)
table_y = Inches(2.5)
col_w = Inches(2.4)
row_h = Inches(0.45)

headers = ['Feature', 'Mealvana', 'Fuelin', 'Hexis', 'MFP/Cronometer']
rows = [
    ['Race-day fueling plans', '\u2713', 'Partial', '\u2717', '\u2717'],
    ['Food preference learning', '\u2713', '\u2717', '\u2717', '\u2717'],
    ['Coach Mode (B2B)', '\u2713', '\u2717', '\u2717', '\u2717'],
    ['TrainingPeaks + Garmin', '\u2713', '\u2713', 'Partial', '\u2717'],
    ['Offline-first', '\u2713', '\u2717', '\u2717', '\u2713'],
    ['ACSM/ISSN evidence base', '\u2713', 'Partial', 'Partial', '\u2717'],
    ['Weather-adjusted hydration', '\u2713', '\u2717', '\u2717', '\u2717'],
    ['Price (monthly)', '$9.99', '~$15', '~$12', 'Free/$10'],
]

# Header row
for j, header in enumerate(headers):
    x = table_x + j * col_w
    card = add_shape(slide, x, table_y, col_w - Inches(0.05), row_h, BLACKBERRY)
    color = ORANGE if header == 'Mealvana' else WHITE
    add_text_box(slide, x + Inches(0.15), table_y + Inches(0.05), col_w - Inches(0.3), row_h,
                 header, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)

# Data rows
for i, row in enumerate(rows):
    y = table_y + (i + 1) * row_h
    bg = WHITE if i % 2 == 0 else RGBColor(0xEE, 0xEC, 0xE6)
    for j, cell in enumerate(row):
        x = table_x + j * col_w
        add_shape(slide, x, y, col_w - Inches(0.05), row_h, bg)

        if cell == '\u2713' and j == 1:  # Mealvana checkmarks
            color = RGBColor(0x22, 0xAA, 0x22)
        elif cell == '\u2713':
            color = RGBColor(0x22, 0xAA, 0x22)
        elif cell == '\u2717':
            color = RGBColor(0xCC, 0x33, 0x33)
        else:
            color = TEXT_DARK

        bold = (j == 1)  # Bold Mealvana column
        add_text_box(slide, x + Inches(0.15), y + Inches(0.05), col_w - Inches(0.3), row_h,
                     cell, font_size=12, color=color, bold=bold,
                     alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

add_slide_number(slide, 10)


# ════════════════════════════════════════════
# SLIDE 11: TEAM
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(3), Inches(0.5),
             'TEAM', font_size=14, color=ORANGE, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.1), Inches(10), Inches(1.0),
             'Data science + engineering + domain expertise.',
             font_size=36, color=WHITE, bold=True)

# Founder cards
founders = [
    ('Xuan Huang', 'CEO & Co-founder',
     'Statistician and data scientist. Built Mealvana\'s personalized nutrition algorithms. Deep expertise in ML, recommendation systems, and statistical modeling. Prior startup experience. Birmingham ecosystem connections via Innovation Depot, Forge, and Alabama Launchpad.',
     '[Photo placeholder]'),
    ('Lee Martin', 'CTO & Co-founder',
     'Full-stack engineer and competitive triathlete. Architected the dual AI + algorithmic system, offline-first architecture, and Coach Mode web portal. Deep Flutter, Supabase, and AI/ML integration expertise. Domain expert as an endurance athlete.',
     '[Photo placeholder]'),
]

for i, (name, role, bio, photo) in enumerate(founders):
    x = Inches(1.2 + i * 5.5)
    y = Inches(2.3)

    card = add_shape(slide, x, y, Inches(5.0), Inches(3.0), BLACKBERRY_LIGHT)

    # Photo placeholder
    photo_shape = add_shape(slide, x + Inches(0.3), y + Inches(0.3), Inches(1.2), Inches(1.2), RGBColor(0x5A, 0x33, 0x66))
    add_text_box(slide, x + Inches(0.35), y + Inches(0.65), Inches(1.1), Inches(0.5),
                 photo, font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(1.7), y + Inches(0.3), Inches(3.0), Inches(0.4),
                 name, font_size=22, color=WHITE, bold=True)

    add_text_box(slide, x + Inches(1.7), y + Inches(0.7), Inches(3.0), Inches(0.3),
                 role, font_size=14, color=ORANGE, bold=True)

    add_text_box(slide, x + Inches(0.3), y + Inches(1.7), Inches(4.4), Inches(1.2),
                 bio, font_size=12, color=CREAM, line_spacing=1.3)

# Advisors / extended team
add_text_box(slide, Inches(1.2), Inches(5.6), Inches(10), Inches(0.4),
             'Extended Team', font_size=16, color=ELECTROLYTE, bold=True)

advisors = 'Dr. Rachel Mitchell (Sports Dietitian, clinical validation)  |  Rui Tian (User Researcher)  |  Haibin Zhuang (Engineering Consultant)'
add_text_box(slide, Inches(1.2), Inches(6.1), Inches(10), Inches(0.4),
             advisors, font_size=13, color=CREAM)

# Next hire
add_text_box(slide, Inches(1.2), Inches(6.6), Inches(10), Inches(0.4),
             'Next Hire:  Head of Growth / Sales & Partnerships \u2014 our most critical gap',
             font_size=14, color=ORANGE, bold=True)

add_slide_number(slide, 11)


# ════════════════════════════════════════════
# SLIDE 12: THE ASK
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(1.2), Inches(0.6), Inches(3), Inches(0.5),
             'THE ASK', font_size=14, color=ORANGE, bold=True)

add_text_box(slide, Inches(1.2), Inches(1.1), Inches(10), Inches(1.0),
             'Pre-seed: $250K\u2013$500K',
             font_size=42, color=BLACKBERRY, bold=True)

add_accent_bar(slide, Inches(1.2), Inches(2.3), Inches(1.0), Pt(3), ORANGE)

add_text_box(slide, Inches(1.2), Inches(2.6), Inches(10), Inches(0.5),
             'First outside capital beyond $25K Alabama Launchpad prize.',
             font_size=18, color=TEXT_DARK)

# Use of funds
add_text_box(slide, Inches(1.2), Inches(3.5), Inches(5), Inches(0.5),
             'Use of Funds', font_size=20, color=BLACKBERRY, bold=True)

funds = [
    ('Head of Growth hire', '40%', 'User acquisition, coach outreach, partnerships'),
    ('Marketing & acquisition', '25%', 'App Store ads, influencer partnerships, event sponsorships'),
    ('Product development', '20%', 'Garmin integration depth, analytics, onboarding optimization'),
    ('Operations & runway', '15%', 'Infrastructure, tools, 18-month runway extension'),
]

for i, (item, pct, detail) in enumerate(funds):
    y = Inches(4.1 + i * 0.7)

    # Percentage bar
    bar_width = float(pct.replace('%', '')) / 100 * 4.5
    add_shape(slide, Inches(1.2), y, Inches(bar_width), Inches(0.35), ORANGE)
    add_shape(slide, Inches(1.2) + Inches(bar_width), y, Inches(4.5 - bar_width), Inches(0.35),
              RGBColor(0xE0, 0xDD, 0xD5))

    add_text_box(slide, Inches(1.3), y + Inches(0.02), Inches(3), Inches(0.3),
                 f'{pct}  {item}', font_size=13, color=WHITE, bold=True)

    add_text_box(slide, Inches(6.0), y + Inches(0.02), Inches(5), Inches(0.3),
                 detail, font_size=13, color=TEXT_DARK)

# Milestones
add_text_box(slide, Inches(7.5), Inches(3.5), Inches(5), Inches(0.5),
             '18-Month Milestones', font_size=20, color=BLACKBERRY, bold=True)

milestones = [
    ('Q2 2026', 'Subscription revenue live'),
    ('Q3 2026', '100+ paying subscribers'),
    ('Q3 2026', '10+ coaches onboarded'),
    ('Q4 2026', '$12K-$70K ARR'),
    ('Q1 2027', 'B2C + B2B PMF proven'),
    ('Q2 2027', 'Series seed readiness'),
]

for i, (date, milestone) in enumerate(milestones):
    y = Inches(4.1 + i * 0.5)
    add_text_box(slide, Inches(7.5), y, Inches(1.2), Inches(0.4),
                 date, font_size=13, color=ORANGE, bold=True)
    add_text_box(slide, Inches(8.8), y, Inches(3.5), Inches(0.4),
                 milestone, font_size=13, color=TEXT_DARK)

add_slide_number(slide, 12)


# ════════════════════════════════════════════
# SLIDE 13: CLOSING / VISION
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)

add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.5),
             'OUR VISION', font_size=14, color=ORANGE, bold=True)

add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(2),
             'Become the nutrition intelligence\nlayer for endurance sports.',
             font_size=44, color=WHITE, bold=True)

add_accent_bar(slide, Inches(1.5), Inches(4.4), Inches(1.2), Pt(4), ORANGE)

add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(1.0),
             'Training platforms tell athletes how to train.\nMealvana tells them how to fuel.',
             font_size=22, color=CREAM)

add_text_box(slide, Inches(1.5), Inches(6.0), Inches(5), Inches(0.4),
             'endurance.mealvana.io', font_size=16, color=ELECTROLYTE, bold=False)

add_text_box(slide, Inches(1.5), Inches(6.4), Inches(5), Inches(0.4),
             'xh.analytics@gmail.com', font_size=14, color=LIGHT_GRAY)


# ════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════
output_path = os.path.expanduser('~/development/mealvana_endurance/docs/_archived/business/Mealvana_Endurance_Pitch_Deck.pptx')
prs.save(output_path)
print(f'Saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
