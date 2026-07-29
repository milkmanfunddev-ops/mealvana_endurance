#!/usr/bin/env python3
"""
Generate Mealvana Endurance presentation for B.E.S.T. Coaches
March 2, 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors ──────────────────────────────────────────────
BLACKBERRY      = RGBColor(0x3D, 0x1F, 0x47)
BLACKBERRY_LIGHT= RGBColor(0x4A, 0x28, 0x54)
ORANGE          = RGBColor(0xFF, 0x8B, 0x3D)
ORANGE_LIGHT    = RGBColor(0xFF, 0xA0, 0x5A)
TEAL            = RGBColor(0x5D, 0xE4, 0xD3)
TEAL_DARK       = RGBColor(0x3F, 0xD4, 0xC0)
PINK            = RGBColor(0xE8, 0x43, 0x93)
CREAM           = RGBColor(0xF5, 0xF3, 0xED)
WHITE           = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT      = RGBColor(0xF8, 0xF6, 0xEB)
TEXT_MUTED      = RGBColor(0xBB, 0xAA, 0xC3)
DARK_BG         = RGBColor(0x2D, 0x15, 0x35)

# ── Paths ─────────────────────────────────────────────────────
BASE = "/Users/leemartin/development/mealvana_endurance"
LOGO = os.path.join(BASE, "assets/images/endurance_welcome_logo.png")
FOOD_LOG_DIARY = os.path.join(BASE, "docs/food_logging/f_new_2.png")
FOOD_LOG_AI    = os.path.join(BASE, "docs/food_logging/f_new_3.png")
FOOD_LOG_HOME  = os.path.join(BASE, "docs/food_logging/f_new_1.png")
OUTPUT         = os.path.join(BASE, "docs/presentations/BEST_Coaches_Mealvana_March2026.pptx")

# ── Presentation Setup ────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, corner_radius=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_LIGHT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=TEXT_LIGHT, bullet_color=ORANGE, spacing=Pt(8)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.level = 0

        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "  "  # indent
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Calibri"

        # Actual text
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"

    return txBox


def add_accent_line(slide, left, top, width, color=ORANGE):
    """Add a horizontal accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_icon_bullet_list(slide, left, top, width, height, items, font_size=16):
    """Add a list with colored dot bullets. items = [(text, dot_color), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, dot_color) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(10)

        dot = p.add_run()
        dot.text = "\u25CF  "
        dot.font.size = Pt(font_size)
        dot.font.color.rgb = dot_color
        dot.font.name = "Calibri"

        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = TEXT_LIGHT
        run.font.name = "Calibri"

    return txBox


def section_header(slide, text, subtitle=None):
    """Add section header at top-left with accent line."""
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                 text, font_size=36, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.25), Inches(2.5), ORANGE)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.45), Inches(10), Inches(0.6),
                     subtitle, font_size=18, color=TEXT_MUTED)


def add_slide_number(slide, num, total):
    """Add small slide number at bottom-right."""
    add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4),
                 f"{num} / {total}", font_size=10, color=TEXT_MUTED,
                 alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

# Subtle gradient overlay shape
add_shape(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BLACKBERRY, 0)

# Logo
if os.path.exists(LOGO):
    slide.shapes.add_picture(LOGO, Inches(5.5), Inches(0.8), Inches(2.3), Inches(2.3))

# Title text
add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(1.2),
             "MEALVANA ENDURANCE", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Accent line under title
add_accent_line(slide, Inches(5.2), Inches(4.4), Inches(3.0), ORANGE)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(4.7), Inches(10.3), Inches(0.6),
             "Science-Based Nutrition Planning for Endurance Athletes",
             font_size=22, color=TEAL, alignment=PP_ALIGN.CENTER)

# Event info
add_text_box(slide, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.5),
             "B.E.S.T. Coaches Presentation  |  March 2, 2026",
             font_size=16, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


TOTAL_SLIDES = 22

# ══════════════════════════════════════════════════════════════
# SLIDE 2: Agenda
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "Today's Agenda")

agenda_items = [
    ("The Nutrition Gap in Coaching", ORANGE),
    ("What is Mealvana Endurance?", ORANGE),
    ("Key Features & How It Works", ORANGE),
    ("Live App Demo", TEAL),
    ("Coach Mode: Your Dashboard", ORANGE),
    ("Coach Mode Demo", TEAL),
    ("Product Roadmap & Upcoming Features", ORANGE),
    ("Design Previews: Food Logging & More", ORANGE),
    ("Pricing & Pilot Coach Program", ORANGE),
    ("Q&A", TEAL),
]
add_icon_bullet_list(slide, Inches(1.0), Inches(2.0), Inches(5.5), Inches(5.0),
                     agenda_items, font_size=20)

# Right side time estimates
time_items = [
    "5 min", "5 min", "5 min", "10 min", "5 min",
    "5 min", "3 min", "3 min", "3 min", "5 min"
]
for i, t in enumerate(time_items):
    add_text_box(slide, Inches(7.5), Inches(2.05 + i * 0.47), Inches(1.5), Inches(0.4),
                 t, font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.LEFT)

add_slide_number(slide, 2, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 3: The Problem
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "The Nutrition Gap", "Why endurance athletes are under-fueled")

# Stat boxes
stats = [
    ("60%", "of age-group athletes\nwing race nutrition", ORANGE),
    ("30-50%", "performance drop from\npoor fueling strategy", PINK),
    ("2-3 hrs/wk", "coaches spend on\nmanual nutrition planning", TEAL),
]

for i, (big, small, color) in enumerate(stats):
    x = Inches(1.0 + i * 4.0)
    y = Inches(2.2)
    # Card background
    add_shape(slide, x, y, Inches(3.5), Inches(2.8), BLACKBERRY_LIGHT, 0.05)
    # Big number
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.0), Inches(1.0),
                 big, font_size=44, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, x + Inches(0.3), y + Inches(1.4), Inches(3.0), Inches(1.2),
                 small, font_size=16, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Quote
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10.0), Inches(1.0),
             '"What surprised us was how many coaches mentioned the fueling knowledge gap."',
             font_size=18, color=CREAM, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(6.1), Inches(10.0), Inches(0.4),
             "-- Mealvana Coach Research, 2026",
             font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 3, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 4: What is Mealvana Endurance?
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "What is Mealvana Endurance?",
               "A personalized nutritionist in every athlete's pocket")

features = [
    "AI-powered nutrition plans generated in seconds",
    "Science-based: ACSM formulas & evidence-based guidelines",
    "Respects athlete food preferences (3-tier like/willing/dislike system)",
    "Integrates with Training Peaks & Final Surge",
    "Offline-first: works without internet, syncs when connected",
    "Multi-sport: running, cycling, swimming & brick workouts",
    "Phase-specific plans: Before, During & After activity",
]
add_bullet_list(slide, Inches(1.0), Inches(2.0), Inches(6.5), Inches(5.0),
                features, font_size=18, bullet_color=TEAL)

# Right side: logo
if os.path.exists(LOGO):
    slide.shapes.add_picture(LOGO, Inches(9.0), Inches(2.5), Inches(3.0), Inches(3.0))

add_slide_number(slide, 4, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 5: How It Works - Science
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "How It Works", "Evidence-based nutrition science meets AI")

# Three columns
cols = [
    ("BEFORE", ORANGE, [
        "1-4g/kg carbs based on time available",
        "Low fat, low fiber for GI comfort",
        "Timing-aware: 15 min to 4 hrs before",
        "Personalized food recommendations",
    ]),
    ("DURING", TEAL, [
        "30-90g/hr carbs (gut training level)",
        "300-1000mg/hr sodium (heat-adjusted)",
        "0.4-0.8L/hr hydration",
        "Portable, easy-to-consume foods",
    ]),
    ("AFTER", PINK, [
        "1g/kg carbs within 2 hours",
        "0.2g/kg protein for recovery",
        "Rehydration targets",
        "Real food preferences respected",
    ]),
]

for i, (title, color, items) in enumerate(cols):
    x = Inches(0.8 + i * 4.2)
    y = Inches(2.0)
    # Card
    add_shape(slide, x, y, Inches(3.8), Inches(4.5), BLACKBERRY_LIGHT, 0.04)
    # Colored top bar
    add_shape(slide, x, y, Inches(3.8), Pt(4), color)
    # Title
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.5),
                 title, font_size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Items
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.3), y + Inches(1.0 + j * 0.7),
                     Inches(3.2), Inches(0.6),
                     f"\u25B8 {item}", font_size=14, color=TEXT_LIGHT)

add_slide_number(slide, 5, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 6: Key Feature - Smart Nutrition Plans
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "Smart Nutrition Plans", "AI + Linear Programming Optimization")

items = [
    ("Dual-System Reliability", "AI primary system for personalization, algorithmic fallback ensures plans ALWAYS generate (<1 sec)", TEAL),
    ("Multi-Objective Optimization", "Simultaneous optimization of carbs, protein, fat, sodium & hydration using linear programming", ORANGE),
    ("Safety-First Design", "No unsafe food combinations (no oatmeal during runs), GI sensitivity tracking, sodium limits", PINK),
    ("Weather-Aware", "Adjusts hydration and sodium based on temperature, humidity & indoor/outdoor", TEAL),
    ("Gut Training Levels", "Graduated carb targets based on athlete's gut training experience", ORANGE),
]

for i, (title, desc, color) in enumerate(items):
    y = Inches(1.8 + i * 1.05)
    # Colored dot
    dot_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.1), Inches(0.15), Inches(0.15))
    dot_shape.fill.solid()
    dot_shape.fill.fore_color.rgb = color
    dot_shape.line.fill.background()
    # Title
    add_text_box(slide, Inches(1.4), y - Inches(0.05), Inches(4.0), Inches(0.4),
                 title, font_size=18, color=color, bold=True)
    # Description
    add_text_box(slide, Inches(1.4), y + Inches(0.3), Inches(10.5), Inches(0.5),
                 desc, font_size=14, color=TEXT_MUTED)

add_slide_number(slide, 6, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 7: Training Platform Integration
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "Training Platform Integration",
               "Works with the tools coaches already use")

# Two columns
tp_items = [
    "Automatic workout sync from Training Peaks",
    "One-click OAuth connection (no credential sharing)",
    "Smart change detection when schedules update",
    "Automatic nutrition plan staleness flagging",
]
fs_items = [
    "Full Final Surge integration",
    "Pull athlete schedules automatically",
    "Support for all activity types",
    "Bidirectional sync capability",
]

# Training Peaks card
add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), BLACKBERRY_LIGHT, 0.04)
add_text_box(slide, Inches(1.2), Inches(2.3), Inches(4.5), Inches(0.5),
             "Training Peaks", font_size=24, color=ORANGE, bold=True)
add_bullet_list(slide, Inches(1.2), Inches(3.0), Inches(4.8), Inches(3.0),
                tp_items, font_size=16, bullet_color=ORANGE)

# Final Surge card
add_shape(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5), BLACKBERRY_LIGHT, 0.04)
add_text_box(slide, Inches(7.4), Inches(2.3), Inches(4.5), Inches(0.5),
             "Final Surge", font_size=24, color=TEAL, bold=True)
add_bullet_list(slide, Inches(7.4), Inches(3.0), Inches(4.8), Inches(3.0),
                fs_items, font_size=16, bullet_color=TEAL)

# Bottom note
add_text_box(slide, Inches(1.0), Inches(6.7), Inches(11.0), Inches(0.4),
             "More platform integrations coming soon (Strava, Garmin, and more)",
             font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 7, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 8: Food Preferences
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "Food Preferences System",
               "Athletes eat what they actually like")

tiers = [
    ("LIKED", "Strong preference (+20 priority)\nRecommended first in plans", TEAL, "\u2764"),
    ("WILLING TO TRY", "Moderate preference (+5 priority)\nUsed when needed for targets", ORANGE, "\u2B50"),
    ("DISLIKED / ALLERGIC", "Automatically excluded\nNever appears in plans", PINK, "\u26D4"),
]

for i, (title, desc, color, icon) in enumerate(tiers):
    x = Inches(0.8 + i * 4.2)
    y = Inches(2.2)
    add_shape(slide, x, y, Inches(3.8), Inches(2.5), BLACKBERRY_LIGHT, 0.04)
    add_shape(slide, x, y, Inches(3.8), Pt(4), color)
    add_text_box(slide, x, y + Inches(0.3), Inches(3.8), Inches(0.5),
                 f"{icon}  {title}", font_size=20, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.2), Inches(1.2),
                 desc, font_size=15, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Bottom benefits
benefits = [
    "100+ foods with detailed nutrition data and phase appropriateness",
    "Phase-specific intelligence: right foods for right timing",
    "Higher adherence = better race results",
]
add_bullet_list(slide, Inches(1.0), Inches(5.2), Inches(11.0), Inches(2.0),
                benefits, font_size=16, bullet_color=TEAL)

add_slide_number(slide, 8, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 9: Live App Demo
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.0),
             "LIVE APP DEMO", font_size=56, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(5.2), Inches(3.5), Inches(3.0), TEAL)
add_text_box(slide, Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.6),
             "Creating a nutrition plan from scratch",
             font_size=22, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

demo_points = [
    "Create a new activity (marathon training run)",
    "Set distance, pace, and conditions",
    "Generate personalized nutrition plan",
    "Review Before / During / After phases",
    "Swap foods and adjust quantities",
]
add_bullet_list(slide, Inches(3.5), Inches(4.6), Inches(6.0), Inches(2.5),
                demo_points, font_size=18, bullet_color=TEAL)

add_slide_number(slide, 9, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 10: Coach Mode Overview
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "Coach Mode", "Manage your athletes' nutrition from one place")

features_left = [
    "Many-to-many: manage multiple athletes",
    "Athlete codes for easy pairing (ATH-XXXXXXXX)",
    "Permission-based access levels",
    "Real-time messaging with athletes",
    "Comment on specific activities or plans",
]
features_right = [
    "View athlete nutrition plans (read or edit)",
    "Track completion rates & macro adherence",
    "Coach directory for athlete discovery",
    "Seamless coach/athlete mode toggle",
    "Web-based dashboard (coming soon)",
]

# Left column
add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.8), BLACKBERRY_LIGHT, 0.04)
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(4.5), Inches(0.5),
             "Athlete Management", font_size=20, color=ORANGE, bold=True)
add_bullet_list(slide, Inches(1.2), Inches(2.8), Inches(4.8), Inches(3.5),
                features_left, font_size=16, bullet_color=ORANGE)

# Right column
add_shape(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.8), BLACKBERRY_LIGHT, 0.04)
add_text_box(slide, Inches(7.4), Inches(2.2), Inches(4.5), Inches(0.5),
             "Nutrition Oversight", font_size=20, color=TEAL, bold=True)
add_bullet_list(slide, Inches(7.4), Inches(2.8), Inches(4.8), Inches(3.5),
                features_right, font_size=16, bullet_color=TEAL)

add_slide_number(slide, 10, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 11: Coach Mode - Dashboard Detail
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "Coach Dashboard", "Everything at a glance")

# Workflow steps
steps = [
    ("1", "REGISTER", "Submit coach application\nfor verification", ORANGE),
    ("2", "CONNECT", "Invite athletes via\nathlete code (ATH-XXXX)", TEAL),
    ("3", "REVIEW", "View plans, track\nadherence, add comments", ORANGE),
    ("4", "COMMUNICATE", "Direct messaging\nand plan annotations", TEAL),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.6 + i * 3.2)
    y = Inches(2.2)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.1), y, Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.size = Pt(24)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "Calibri"

    # Title
    add_text_box(slide, x, y + Inches(0.9), Inches(3.0), Inches(0.5),
                 title, font_size=20, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, x, y + Inches(1.4), Inches(3.0), Inches(1.0),
                 desc, font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Bottom detail
add_shape(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.0), BLACKBERRY_LIGHT, 0.03)
add_text_box(slide, Inches(1.2), Inches(5.0), Inches(11.0), Inches(0.4),
             "Dual-Write Architecture: Changes sync in real-time across all devices",
             font_size=16, color=TEAL, bold=True)
detail_items = [
    "Coach writes to Supabase immediately (no sync delays)",
    "Athletes see coach feedback and plan changes instantly",
    "Full offline support: log locally, sync when connected",
]
add_bullet_list(slide, Inches(1.2), Inches(5.5), Inches(10.5), Inches(1.5),
                detail_items, font_size=14, bullet_color=TEAL)

add_slide_number(slide, 11, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 12: Coach Mode - Messaging
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "Coach Communication", "Stay connected with your athletes")

msg_types = [
    ("General Chat", "Direct messaging for\ngeneral coaching conversation", ORANGE),
    ("Activity Comments", "Attach feedback to specific\nworkouts and training sessions", TEAL),
    ("Nutrition Plan Notes", "Comment on nutrition plans\nwith specific recommendations", PINK),
]

for i, (title, desc, color) in enumerate(msg_types):
    x = Inches(0.8 + i * 4.2)
    y = Inches(2.2)
    add_shape(slide, x, y, Inches(3.8), Inches(2.2), BLACKBERRY_LIGHT, 0.04)
    add_shape(slide, x, y, Inches(3.8), Pt(4), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.5),
                 title, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.2), Inches(1.0),
                 desc, font_size=15, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Additional features
extras = [
    "Real-time message delivery via Supabase Realtime subscriptions",
    "Read tracking and unread message counts",
    "Cross-device sync across all athlete devices",
    "Thread-based conversations organized by context",
]
add_bullet_list(slide, Inches(1.0), Inches(5.0), Inches(11.0), Inches(2.0),
                extras, font_size=16, bullet_color=ORANGE)

add_slide_number(slide, 12, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 13: Coach Mode Demo
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.0),
             "COACH MODE DEMO", font_size=56, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(5.2), Inches(3.5), Inches(3.0), ORANGE)
add_text_box(slide, Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.6),
             "Managing athletes as a coach",
             font_size=22, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

demo_points = [
    "Coach dashboard overview",
    "Invite an athlete using athlete code",
    "View athlete's activity and nutrition plan",
    "Add comments and feedback on a plan",
    "Messaging between coach and athlete",
]
add_bullet_list(slide, Inches(3.5), Inches(4.6), Inches(6.0), Inches(2.5),
                demo_points, font_size=18, bullet_color=ORANGE)

add_slide_number(slide, 13, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 14: Product Roadmap
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "Product Roadmap", "What's coming next")

# Roadmap items from Notion - organized by status
roadmap_sections = [
    ("IN PROGRESS", ORANGE, [
        ("Coach / Dietitian Dashboard", "Pro", "Ship: Feb 27", "Developing"),
        ("Better Food Recommendations", "Free", "Ship: Feb 23", "Developing"),
        ("Critical Usability Fixes", "Free", "Ship: Mar 6", "Design"),
        ("Race Day Spreadsheet", "Pro", "Ship: Feb 27", "Design"),
        ("Meal Planning", "Cookies*", "Ship: Apr 17", "Design"),
    ]),
    ("COMING SOON", TEAL, [
        ("Personal Fueling Templates", "Pro", "Ship: Mar 13", "Not started"),
        ("Payment Infrastructure", "Free", "Ship: Mar 13", "Not started"),
        ("Other Platform Integrations", "Pro", "Ship: Mar 27", "Not started"),
        ("Gamification & Rewards", "Free", "Ship: Apr 3", "Not started"),
        ("Fueling Knowledge Gap Content", "Pro", "Ship: Apr 24", "Not started"),
    ]),
]

for col, (section_title, color, items) in enumerate(roadmap_sections):
    x = Inches(0.6 + col * 6.4)
    y_start = Inches(2.0)

    add_text_box(slide, x, y_start, Inches(5.8), Inches(0.5),
                 section_title, font_size=18, color=color, bold=True)
    add_accent_line(slide, x, y_start + Inches(0.45), Inches(2.0), color)

    for i, (name, tier, ship, status) in enumerate(items):
        y = y_start + Inches(0.7 + i * 0.85)
        add_shape(slide, x, y, Inches(5.8), Inches(0.7), BLACKBERRY_LIGHT, 0.03)

        # Name
        add_text_box(slide, x + Inches(0.2), y + Inches(0.05), Inches(3.5), Inches(0.35),
                     name, font_size=14, color=WHITE, bold=True)
        # Tier badge
        tier_color = ORANGE if tier == "Pro" else TEAL if tier == "Free" else PINK
        add_text_box(slide, x + Inches(3.8), y + Inches(0.05), Inches(0.8), Inches(0.35),
                     tier, font_size=11, color=tier_color, bold=True)
        # Ship date
        add_text_box(slide, x + Inches(4.6), y + Inches(0.05), Inches(1.1), Inches(0.35),
                     ship, font_size=11, color=TEXT_MUTED)
        # Status
        add_text_box(slide, x + Inches(0.2), y + Inches(0.35), Inches(5.0), Inches(0.3),
                     status, font_size=11, color=TEXT_MUTED)

# Footer note
add_text_box(slide, Inches(0.6), Inches(6.8), Inches(12.0), Inches(0.4),
             "*Cookies = token-based premium feature  |  Timeline subject to change",
             font_size=12, color=TEXT_MUTED)

add_slide_number(slide, 14, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 15: Food Logging - Design Preview
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "Coming Soon: Food Logging",
               "AI-powered daily nutrition tracking")

# Add food logging screenshots
screenshots = [
    (FOOD_LOG_HOME, "Activities View"),
    (FOOD_LOG_DIARY, "Nutrition Diary"),
    (FOOD_LOG_AI, "AI Food Logger"),
]

for i, (path, label) in enumerate(screenshots):
    x = Inches(0.8 + i * 4.2)
    if os.path.exists(path):
        # Add screenshot with constrained height
        slide.shapes.add_picture(path, x + Inches(0.4), Inches(2.2), Inches(3.0), Inches(4.5))
    # Label
    add_text_box(slide, x, Inches(6.8), Inches(3.8), Inches(0.4),
                 label, font_size=14, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 15, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 16: Food Logging Features
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "Food Logging Features",
               "Multiple ways to track nutrition")

features_list = [
    ("AI Text Parsing", 'Type "grilled chicken 6oz with rice" and AI parses it into structured food items', ORANGE),
    ("AI Photo Recognition", "Take a photo of your plate and AI identifies foods and estimates portions", TEAL),
    ("Barcode Scanning", "Scan product barcodes for instant nutrition data", ORANGE),
    ("Database Search", "Search 100+ foods from our curated database", TEAL),
    ("Daily Macro Dashboard", "Track Calories, Carbs, Protein, Fat, Sodium & Hydration at a glance", ORANGE),
    ("Timeline Integration", "Meals mixed chronologically with activities for a complete daily view", TEAL),
]

for i, (title, desc, color) in enumerate(features_list):
    y = Inches(1.8 + i * 0.88)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.08), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text_box(slide, Inches(1.4), y - Inches(0.08), Inches(4.0), Inches(0.35),
                 title, font_size=17, color=color, bold=True)
    add_text_box(slide, Inches(1.4), y + Inches(0.25), Inches(10.5), Inches(0.45),
                 desc, font_size=14, color=TEXT_MUTED)

add_slide_number(slide, 16, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 17: Race Day Spreadsheet
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "Coming Soon: Race Day Spreadsheet",
               "Replace manual spreadsheets with smart race nutrition")

items = [
    "Coaches currently use spreadsheets to plan race day nutrition",
    "Mealvana replicates this function with better convenience and accuracy",
    "Automatically calculated based on race distance, pace, and conditions",
    "Easy to share with athletes through the coach portal",
    "Design currently in progress - prototype available",
]
add_bullet_list(slide, Inches(1.0), Inches(2.2), Inches(11.0), Inches(3.5),
                items, font_size=18, bullet_color=ORANGE)

# Prototype link reference
add_shape(slide, Inches(3.0), Inches(5.5), Inches(7.3), Inches(1.2), BLACKBERRY_LIGHT, 0.04)
add_text_box(slide, Inches(3.3), Inches(5.7), Inches(6.7), Inches(0.4),
             "Figma Prototype Available", font_size=20, color=TEAL, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.3), Inches(6.15), Inches(6.7), Inches(0.4),
             "We'd love your feedback on the race day nutrition design!",
             font_size=14, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 17, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 18: Blank Design Slide 1
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "Design Preview", "Additional designs")

# Placeholder text
add_text_box(slide, Inches(2.0), Inches(3.0), Inches(9.3), Inches(1.5),
             "[Add your Figma designs or screenshots here]",
             font_size=28, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Dashed border placeholder
add_shape(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.5), BLACKBERRY_LIGHT, 0.03)

add_slide_number(slide, 18, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 19: Blank Design Slide 2
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "Design Preview", "Additional designs")

add_text_box(slide, Inches(2.0), Inches(3.0), Inches(9.3), Inches(1.5),
             "[Add your Figma designs or screenshots here]",
             font_size=28, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.5), BLACKBERRY_LIGHT, 0.03)

add_slide_number(slide, 19, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 20: Payment Strategy
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BLACKBERRY)
section_header(slide, "Payment Strategy", "Simple, fair pricing for athletes and coaches")

# Athlete tier
add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), BLACKBERRY_LIGHT, 0.04)
add_text_box(slide, Inches(1.2), Inches(2.3), Inches(4.7), Inches(0.5),
             "ATHLETE SUBSCRIPTION", font_size=18, color=ORANGE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.2), Inches(3.0), Inches(4.7), Inches(1.0),
             "$19.99", font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.2), Inches(3.9), Inches(4.7), Inches(0.4),
             "per month", font_size=18, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
athlete_features = [
    "All nutrition plan features",
    "Training Peaks & Final Surge sync",
    "Food preference personalization",
    "Offline-first mobile app",
]
add_bullet_list(slide, Inches(1.5), Inches(4.5), Inches(4.0), Inches(2.0),
                athlete_features, font_size=14, bullet_color=ORANGE)

# Coach tier
add_shape(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5), BLACKBERRY_LIGHT, 0.04)
add_text_box(slide, Inches(7.4), Inches(2.3), Inches(4.7), Inches(0.5),
             "COACH SUBSCRIPTION", font_size=18, color=TEAL, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(7.4), Inches(3.0), Inches(4.7), Inches(1.0),
             "$20 - $50", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(7.4), Inches(3.9), Inches(4.7), Inches(0.4),
             "per month (based on athlete count)", font_size=16, color=TEXT_MUTED,
             alignment=PP_ALIGN.CENTER)
coach_features = [
    "Starter: $20/mo (up to 20 athletes)",
    "Pro: $50/mo (unlimited athletes)",
    "Full dashboard & analytics",
    "Messaging & plan annotations",
]
add_bullet_list(slide, Inches(7.7), Inches(4.5), Inches(4.0), Inches(2.0),
                coach_features, font_size=14, bullet_color=TEAL)

# Note about payment infrastructure
add_text_box(slide, Inches(1.0), Inches(6.8), Inches(11.3), Inches(0.4),
             "Exploring options outside Apple App Store for more flexible coach billing (Stripe)",
             font_size=13, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 20, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 21: Pilot Coaches Program
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.8),
             "JOIN OUR PILOT COACHES PROGRAM", font_size=40, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(4.5), Inches(1.85), Inches(4.3), ORANGE)

# Big discount badge
badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.1), Inches(2.3), Inches(3.1), Inches(3.1))
badge.fill.solid()
badge.fill.fore_color.rgb = ORANGE
badge.line.fill.background()
# Text inside badge
add_text_box(slide, Inches(5.3), Inches(2.8), Inches(2.7), Inches(0.8),
             "30%", font_size=64, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.3), Inches(3.6), Inches(2.7), Inches(0.5),
             "OFF", font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# What you get
add_text_box(slide, Inches(1.0), Inches(5.6), Inches(11.3), Inches(0.5),
             "What Pilot Coaches Get:", font_size=20, color=TEAL, bold=True,
             alignment=PP_ALIGN.CENTER)

perks = [
    "30% discount on coach subscription for the first year",
    "Early access to new features before public release",
    "Direct line to the development team for feedback",
    "Shape the product roadmap with your coaching expertise",
    "Priority support and onboarding assistance",
]
add_bullet_list(slide, Inches(3.0), Inches(6.1), Inches(7.5), Inches(2.5),
                perks, font_size=15, bullet_color=ORANGE, spacing=Pt(4))

add_slide_number(slide, 21, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════
# SLIDE 22: Thank You / Q&A
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

# Logo
if os.path.exists(LOGO):
    slide.shapes.add_picture(LOGO, Inches(5.5), Inches(1.0), Inches(2.3), Inches(2.3))

add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.0),
             "Thank You!", font_size=52, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.2), Inches(4.5), Inches(3.0), TEAL)

add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.6),
             "Questions & Discussion", font_size=24, color=TEAL,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.5),
             "MEALVANA ENDURANCE", font_size=18, color=ORANGE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.4),
             "Science-Based Nutrition for Endurance Athletes",
             font_size=16, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 22, TOTAL_SLIDES)


# ── Save ──────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
print(f"\nPresentation saved to:\n{OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
